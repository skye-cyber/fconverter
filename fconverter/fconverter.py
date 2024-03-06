#!/usr/bin/env python3.11.7
import argparse
import os
import pdfminer.high_level
import re
import traceback
import subprocess
from docx import Document
from pptx import Presentation
from .xls2Sql import convert_xlsx_to_database
from .image import enhance_image
from .xlsx import convert_xls_to_word, convert_xls_to_text
from .OCR import ocr_text_extraction
from .xlsx2csv import convert_xlsx_to_csv
from .interactive import interact
from .word_to_pdf import get_word2pdf_files
from .banner import banner
from .mp3 import get_2mp3_files
from .draft import progress_show
import logging
import logging.handlers

logging.basicConfig(level=logging.INFO, format='%(levelname)-8s %(message)s')
logger = logging.getLogger(__name__)


def check_if_running_inside_termux():
    try:
        termux_cmd = ['which', 'apt']
        return subprocess.call(termux_cmd) == 0
    except Exception:
        logger.error('Unable to identify environment!')
        return False


def pdf_to_word(pdf_file, word_file):
    if os.name == 'posix' or os.name == 'ms-dos':
        try:
            from pdf2docx import parse
        except ImportError:
            logger.error('pdf2docx not found \033[34mInstalling...\033[0m')
            subprocess.run(['pip', 'install', 'pdf2docx'])
            pdf_to_txt(pdf_file, word_file)
        try:
            parse(pdf_file, word_file, start=0, end=None)
            return True  # return true if the conversion is successful
            logger.info(f"\033[1;95m Successfully converted{pdf_file} to \
{word_file}\033[0m")
        except Exception as e:
            logger.error(f"Error converting {pdf_file} to {word_file}: {e}")
            with open("conversion.log", "a") as log_file:
                log_file.write(f"Error converting {pdf_file} to {word_file}:\
{e}\n")
            return False  # Return False if conversion fails for all encodings
    elif check_if_running_inside_termux():
        try:
            from PyPDF2 import PdfFileReader
            # Open the pdf file using PyPDF2
            logger.info('Reading pdf..')
            with open(pdf_file, 'rb') as fh:
                reader = PdfFileReader(fh)
                # create a blanks document using python-docx
                doc = Document
                # Add each page content into the Docx document
                for i in range(reader.numPages):
                    page = reader.getPage(i)
                    progress_show(reader.getPageNumber, reader.getNumPages)
                    logger.info(f'\033[32mPage{i}/{reader.numPages}\033[0m')
                    text = page.extractText()  # .replace('\r\n', '\n').strip()
                    # para = doc.add_paragraph(text)
                    doc.add_paragraph(text)
                # Save the resulting File
                doc.save(word_file)
                logger.info(f"\033[1;95m Successfully converted{pdf_file} to \
{word_file}\033[0m")
        except ImportError:
            logger.info('Failed to import PyPDF2 \033[34mInstalling it\033[0m')
            subprocess.run(['pip', 'install', 'PyPDF2'])
            pdf_to_txt(pdf_file, word_file)
        except Exception as e:
            logger.info(f'\033[31m{e}\033[0m')


def word_to_ppt(word_file, ppt_file):
    x = 0
    for x in range[:1000]:
        x += x
    try:
        document = Document(word_file)
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[1]
        for paragraph in document.paragraphs:
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = x
            content = slide.placeholders[1]
            content.text = paragraph.text
            presentation.save(ppt_file)
            logger.info(f"\033[1;95mSuccessfully converted {word_file} to \
{ppt_file}\033[0m")
    except Exception as e:
        logger.error(f"Oops something went awry while attempting to convert \
{word_file} to {ppt_file}: {e}")
        traceback_info = traceback.format_exe()
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {word_file} to {ppt_file}:\n\
{traceback_info}\n")


def word_to_txt(word_file, txt_file):
    try:
        doc = Document(word_file)
        with open(txt_file, 'w', encoding='utf-8') as f:
            for paragraph in doc.paragraphs:
                f.write(paragraph.text + '\n')
                logger.info("Processing...")
            logger.info(f"\033[1;95mSuccessfully converted {word_file} to \
{txt_file}\033[0m")
    except Exception as e:
        logger.error(f"Oops something went amiss while attempting to convert \
{word_file} to {txt_file}: {e}")
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {word_file} to {txt_file}:{e}\n")


def pdf_to_txt(pdf_file, txt_file):
    try:
        with open(pdf_file, 'rb') as f:
            text = pdfminer.high_level.extract_text(pdf_file)
        with open(txt_file, 'w', encoding='utf-8') as f:
            f.write(text)
            logger.info(f"\033[1;95mSuccessfully converted {pdf_file} to \
{txt_file}\033[0m")
    except Exception as e:
        logger.error(f"Oops when astray while converting {pdf_file} to \
{txt_file}: {e}")
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {pdf_file} to {txt_file}: {e}\n")


def ppt_to_word(ppt_file, word_file):
    try:
        presentation = Presentation(ppt_file)
        document = Document()

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        new_paragraph = document.add_paragraph()
                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run(run.text)
                            # Preserve bold formatting
                            new_run.bold = run.font.bold
                            # Preserve italic formatting
                            new_run.italic = run.font.italic
                            # Preserve underline formatting
                            new_run.underline = run.font.underline
                            # Preserve font name
                            new_run.font.name = run.font.name
                            # Preserve font size
                            new_run.font.size = run.font.size
                            try:
                                # Preserve font color
                                new_run.font.color.rgb = run.font.color.rgb
                            except AttributeError:
                                # Ignore error and continue without
                                # setting the font color
                                pass
                    # Add a new paragraph after each slide
                    document.add_paragraph()
        document.save(word_file)
        logger.info(f"\033[1;95mSuccessfully converted {ppt_file} to \
{word_file}\033[0m")
    except Exception as e:
        logger.error(f"Oops somethin gwent awry while attempting to convert \
{ppt_file} to {word_file}:\n>>>{e}")
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Oops something went astray while attempting \
convert {ppt_file} to {word_file}:{e}\n")


def text_to_word(text_file, word_file):
    try:
        # Read the text file
        with open(text_file, 'r', encoding='utf-8', errors='ignore') as file:
            text_content = file.read()

        # Filter out non-XML characters
        filtered_content = re.sub(
            r'[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD]+', '', text_content)

        # Create a new Word document
        doc = Document()
        # Add the filtered text content to the document
        doc.add_paragraph(filtered_content)

        # Save the document as a Word file
        doc.save(word_file)
        logger.info(f"\033[1;95mSuccessfully converted {text_file} to \
{word_file}\033[0m")
    except FileExistsError as e:
        logger.error(f"{str(e)}")
    except Exception as e:
        logger.error(f"Oops Unable to perfom requested conversion: {e}\n")
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {text_file} to {word_file}: \
{e}\n")


def main():
    parser = argparse.ArgumentParser(description='''Convert files between
                                                 different formats.''')
    banner
    parser.add_argument('conversion_type', type=int, help='''The type of
conversion to perform (\033[1;96m 1-14 \033[0m).
\033[1;96m 1:\033[0m Word to PDF,  \033[1;96m 2:\033[m PDF to Word,\
\033[1;96m 3:\033[0m Word to PPT, \033[1;96m 4:\033[0m Word to TXT,\
\033[1;96m 5:\033[0m PDF to TXT,  \033[1;96m 6:\033[0mPPTX to Word,\
\033[1;96m 7:\033[0m TXT to Word,   \033[1;96m 8:\033[0m PDF/DOCX/TXT to mp3,\
\033[1;96m 9:\033[0m Image Enhancement \033[1;96m 10:\033[0m XLSX to Word,\
\033[1;96m 11:\033[0m XLSX to TXT \033[1;96m 12:\033[0m Image Text Extraction,\
\033[1;96m 13:\033[0m Convert xlsx to sqlite db  \033[1;96m 14:\033[0m XLSX to\
 CSV\033[1;96m 15:\033[0m Other conversions including the above\n\n''')

    parser.add_argument('input_file', type=str, help='Name of input file')
    parser.add_argument('output_file', type=str, help='Name of output file')
    args = parser.parse_args()
    input_file = args.input_file
    output_file = args.output_file

    if args.conversion_type == 1:
        get_word2pdf_files(input_file, output_file)
    elif args.conversion_type == 2:
        pdf_to_word(input_file, output_file)
    elif args.conversion_type == 3:
        word_to_ppt(input_file, output_file)
    elif args.conversion_type == 4:
        word_to_txt(input_file, output_file)
    elif args.conversion_type == 5:
        pdf_to_txt(input_file, output_file)
    elif args.conversion_type == 6:
        ppt_to_word(input_file, output_file)
    elif args.conversion_type == 7:
        text_to_word(input_file, output_file)
    elif args.conversion_type == 8:
        get_2mp3_files(input_file, output_file)
    elif args.conversion_type == 9:
        enhance_image(input_file, output_file)
    elif args.conversion_type == 10:
        convert_xls_to_word(input_file, output_file)
    elif args.conversion_type == 11:
        convert_xls_to_text(input_file, output_file)
    elif args.conversion_type == 12:
        try:
            from . import OCRbanner
            OCRbanner
        except Exception as e:
            logger.error(f"{e}")
        ocr_text_extraction(input_file, output_file)
    elif args.conversion_type == 13:
        parser.add_argument('table_name', type=str, help='Name of table, name \
for the db')
        convert_xlsx_to_database(input_file, output_file, args.table_name)
    elif args.conversion_type == 14:
        convert_xlsx_to_csv(input_file, output_file)
    elif args.conversion_type == 15:
        interact
    else:
        logger.error("\033[1;91m Invalid conversion type. Please enter a \
number from 1 to 14.\033[0m")


if __name__ == '__main__':
    main()
