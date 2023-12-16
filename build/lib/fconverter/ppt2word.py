from docx import Document
from pptx import Presentation


def ppt_to_word(ppt_file, word_file):
    try:
        presentation = Presentation(ppt_file)
        document = Document()

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        new_paragraph = document.add_paragraph()
                        for run in paragraph.runs:
                            new_run = new_paragraph.add_run(run.text)
                            # Preserve bold formatting
                            new_run.bold = run.font.bold
                            # Preserve italic formatting
                            new_run.italic = run.font.italic
                            # Preserve underline formatting
                            new_run.underline = run.font.underline
                            # Preserve font name
                            new_run.font.name = run.font.name
                            # Preserve font size
                            new_run.font.size = run.font.size
                            try:
                                # Preserve font color
                                new_run.font.color.rgb = run.font.color.rgb
                            except AttributeError:
                                # Ignore error and continue without
                                # setting the font color
                                pass
                    # Add a new paragraph after each slide
                    document.add_paragraph()
        document.save(word_file)
        print(f"\033[1;95mSuccessfully converted {ppt_file} to {word_file}\033[0m")
    except Exception as e:
        print(f"Error converting {ppt_file} to {word_file}:\n>>> {e}")
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {ppt_file} to {word_file}:{e}\n")
