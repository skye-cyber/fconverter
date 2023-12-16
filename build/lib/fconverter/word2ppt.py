from docx import Document
from pptx import Presentation


def word_to_ppt(word_file, ppt_file):
    x = 0
    for x in range[:1000]:
        x += x
    try:
        document = Document(word_file)
        presentation = Presentation()
        slide_layout = presentation.slide_layouts[1]
        for paragraph in document.paragraphs:
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = x
            content = slide.placeholders[1]
            content.text = paragraph.text
            presentation.save(ppt_file)
            print(f"\033[1;95mSuccessfully converted {word_file} to {ppt_file}\033[0m")
    except Exception as e:
        print(f"Error converting {word_file} to {ppt_file}: {e}")
        traceback_info = traceback.format_exe()
        with open("conversion.log", "a") as log_file:
            log_file.write(f"Error converting {word_file} to {ppt_file}:\n{traceback_info}\n")
